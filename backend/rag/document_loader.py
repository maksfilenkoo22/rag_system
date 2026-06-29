from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pypdf import PdfReader


def load_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def load_html(path: Path) -> str:
    html = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    return soup.get_text(separator="\n")


def load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages_text = []

    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""

        if text.strip():
            pages_text.append(f"\n[Страница {page_number}]\n{text}")

    return "\n".join(pages_text)


def load_docx(path: Path) -> str:
    document = Document(str(path))
    parts = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            row_values = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_values.append(cell_text)

            if row_values:
                parts.append(" | ".join(row_values))

    return "\n".join(parts)


def load_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    slides_text = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts = [f"[Слайд {slide_number}]"]

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)

        slides_text.append("\n".join(slide_parts))

    return "\n\n".join(slides_text)


def load_document(path: Path) -> str:
    suffix = path.suffix.lower()

    if suffix == ".txt":
        return load_txt(path)

    if suffix in [".html", ".htm"]:
        return load_html(path)

    if suffix == ".pdf":
        return load_pdf(path)

    if suffix == ".docx":
        return load_docx(path)

    if suffix == ".pptx":
        return load_pptx(path)

    raise ValueError(
        f"Формат файла {suffix} пока не поддерживается. "
        f"Используй pdf, docx, pptx, html или txt."
    )